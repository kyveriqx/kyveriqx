"""
Build a board-of-directors PowerPoint deck for FY 2024-25.
Reads .tmp/financial_summary.json and writes BOD_Presentation_FY24-25.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree
import json, os, base64, io

from config_loader import BRANDING, OUTLOOK, WORKDIR

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TMP = WORKDIR
OUT = os.path.join(WORKDIR, 'BOD_Presentation.pptx')

# Resolve company-level branding
COMPANY_NAME    = BRANDING.get('companyName') or '[COMPANY NAME]'
TAGLINE         = BRANDING.get('tagline') or '[Tagline / Vision Statement]'
VISION          = BRANDING.get('vision') or ''
PREPARED_FOR    = BRANDING.get('preparedFor') or 'Board of Directors'
REPORTING_PERIOD = BRANDING.get('reportingPeriod') or 'FY 2024-25'

def _hex_to_rgb(hexstr, fallback=(0x1F, 0x38, 0x64)):
    if not hexstr:
        return fallback
    h = str(hexstr).lstrip('#')
    if len(h) != 6:
        return fallback
    try:
        return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        return fallback

_PRIMARY = _hex_to_rgb(BRANDING.get('primaryColor'), (0x1F, 0x38, 0x64))
_ACCENT  = _hex_to_rgb(BRANDING.get('accentColor'),  (0xBF, 0x8F, 0x00))

# Decode logo data URL once (if provided)
_LOGO_BYTES = None
_logo_url = BRANDING.get('logoDataUrl')
if _logo_url and isinstance(_logo_url, str) and _logo_url.startswith('data:image'):
    try:
        b64 = _logo_url.split(',', 1)[1]
        _LOGO_BYTES = base64.b64decode(b64)
    except Exception as e:
        print(f"[deck] Failed to decode logo data URL: {e}")

with open(os.path.join(TMP, 'financial_summary.json')) as f:
    s = json.load(f)

pl  = s['pl_actual']
pla = s['pl_annualized']

# ----- Brand palette (driven by branding) -----
NAVY    = RGBColor(*_PRIMARY)
NAVY2   = RGBColor(min(_PRIMARY[0] + 0x10, 0xFF),
                   min(_PRIMARY[1] + 0x20, 0xFF),
                   min(_PRIMARY[2] + 0x30, 0xFF))
GOLD    = RGBColor(*_ACCENT)
GOLD_L  = RGBColor(min(_ACCENT[0] + 0x40, 0xFF),
                   min(_ACCENT[1] + 0x40, 0xFF),
                   min(_ACCENT[2] + 0x60, 0xFF))
LIGHT_N = RGBColor(0xD9, 0xE1, 0xF2)
BG_LIGHT= RGBColor(0xF2, 0xF2, 0xF2)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x26, 0x26, 0x26)
GRAY    = RGBColor(0x59, 0x59, 0x59)
GRAY_L  = RGBColor(0xBF, 0xBF, 0xBF)
GREEN   = RGBColor(0x54, 0x82, 0x35)
RED     = RGBColor(0xC0, 0x00, 0x00)

# ----- Setup 16:9 deck -----
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

def to_cr(v): return v / 1e7
def fmt_cr(v, signed=False):
    if v < 0: return f"(₹{abs(v)/1e7:,.2f} Cr)"
    return f"₹{v/1e7:,.2f} Cr"
def fmt_pct(v, signed=False):
    if v < 0: return f"({abs(v):.1f}%)"
    return f"{v:.1f}%"

def blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, size=14, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, font='Calibri'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        f = run.font
        f.name = font; f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = color
    return tb

def add_title_bar(slide, title, subtitle=None):
    # Brand strip across top
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), NAVY)
    # Title
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.6),
             title, size=26, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.75), Inches(12.3), Inches(0.35),
                 subtitle, size=12, italic=True, color=GRAY)
    # Divider line
    add_rect(slide, Inches(0.5), Inches(1.13), Inches(12.3), Emu(12700), GOLD)

def add_footer(slide, page_no=None):
    add_rect(slide, 0, SLIDE_H - Inches(0.4), SLIDE_W, Inches(0.4), NAVY)
    add_text(slide, Inches(0.3), SLIDE_H - Inches(0.36), Inches(8), Inches(0.32),
             f"{COMPANY_NAME}  |  {PREPARED_FOR}  |  {REPORTING_PERIOD} Performance Review",
             size=10, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if page_no is not None:
        add_text(slide, SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.36), Inches(1.2), Inches(0.32),
                 f"Page {page_no}", size=10, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

# =====================================================
# SLIDE 1 - TITLE
# =====================================================
slide = blank_slide()
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
# accent stripe
add_rect(slide, 0, Inches(4.5), SLIDE_W, Inches(0.08), GOLD)
# Diagonal accent block
shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, SLIDE_W - Inches(4), 0, Inches(4), Inches(3.5))
shp.fill.solid(); shp.fill.fore_color.rgb = NAVY2
shp.line.fill.background()
shp.rotation = 180

add_text(slide, Inches(0.7), Inches(1.2), Inches(11), Inches(0.5),
         "ANNUAL BOARD REVIEW", size=14, bold=True, color=GOLD, font='Calibri')

# Optional logo (top-right)
if _LOGO_BYTES:
    try:
        slide.shapes.add_picture(io.BytesIO(_LOGO_BYTES),
                                 SLIDE_W - Inches(1.6), Inches(0.4),
                                 width=Inches(1.2), height=Inches(1.2))
    except Exception as e:
        print(f"[deck] Logo embed failed: {e}")

add_text(slide, Inches(0.7), Inches(1.8), Inches(11), Inches(1.5),
         COMPANY_NAME, size=54, bold=True, color=WHITE)

add_text(slide, Inches(0.7), Inches(3.0), Inches(11), Inches(0.6),
         TAGLINE, size=18, italic=True, color=LIGHT_N)

add_text(slide, Inches(0.7), Inches(4.9), Inches(11), Inches(0.7),
         "Financial Performance & Strategic Outlook", size=24, bold=True, color=WHITE)

add_text(slide, Inches(0.7), Inches(5.7), Inches(11), Inches(0.5),
         f"{REPORTING_PERIOD}  |  Presented to the {PREPARED_FOR}", size=16, color=GOLD_L)

add_text(slide, Inches(0.7), Inches(6.6), Inches(11), Inches(0.4),
         f"Period covered: {s['period']}", size=11, italic=True, color=LIGHT_N)

# =====================================================
# SLIDE 2 - AGENDA
# =====================================================
slide = blank_slide()
add_title_bar(slide, "Agenda", "What we'll cover today")
add_footer(slide, 2)

agenda = [
    ("01", "Company Snapshot", "Who we are, what we do"),
    ("02", f"Financial Highlights — {REPORTING_PERIOD}", "Headline P&L and margins"),
    ("03", "Profit & Loss Statement", "Detailed line-by-line walkthrough"),
    ("04", "Revenue Trend & Momentum", "Monthly trajectory and seasonality"),
    ("05", "Customer & Vendor Concentration", "Where the money comes from and goes"),
    ("06", "Geographic & Currency Mix", "Domestic vs export footprint"),
    ("07", "Working Capital & Liquidity", "AR, AP, and cash conversion"),
    ("08", "Strategic Outlook FY 2025-26", "Plans, priorities, and growth bets"),
    ("09", "Risks, Asks & Discussion", "Board input required"),
]
y = Inches(1.4)
for num, title, sub in agenda:
    # Number badge
    add_rect(slide, Inches(0.7), y, Inches(0.8), Inches(0.55), NAVY)
    add_text(slide, Inches(0.7), y, Inches(0.8), Inches(0.55), num, size=20, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(1.7), y, Inches(7), Inches(0.3), title, size=16, bold=True, color=NAVY)
    add_text(slide, Inches(1.7), y + Inches(0.28), Inches(11), Inches(0.3), sub,
             size=11, italic=True, color=GRAY)
    y += Inches(0.6)

# =====================================================
# SLIDE 3 - COMPANY SNAPSHOT
# =====================================================
slide = blank_slide()
add_title_bar(slide, "Company Snapshot", f"{COMPANY_NAME} at a glance")
add_footer(slide, 3)

# Left column - About
add_rect(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(5.5), BG_LIGHT)
add_text(slide, Inches(0.7), Inches(1.55), Inches(5.6), Inches(0.4),
         "About Us", size=18, bold=True, color=NAVY)
# Bullets are computed from the user's data + branding — no hardcoded
# sample text. Bullets with no real signal are skipped so the box doesn't
# show "Top customer contributes 0%".
_about_lines = []
if VISION:
    _about_lines.append(f"• {VISION}")
if TAGLINE and TAGLINE != VISION:
    _about_lines.append(f"• {TAGLINE}")
_about_lines.append(f"• Reporting period: {REPORTING_PERIOD}")
_about_lines.append(f"• Prepared for: {PREPARED_FOR}")
_top_c_pre = s.get('top_customers') or []
if _top_c_pre and _top_c_pre[0][1] > 0:
    _rev_sum = sum(a for _, a in _top_c_pre[:10]) or 1
    _share = _top_c_pre[0][1] / _rev_sum * 100
    _about_lines.append(
        f"• Top customer: {str(_top_c_pre[0][0])[:36]} "
        f"({_share:.1f}% of top-10 sales)"
    )
_top_v_pre = s.get('top_vendors') or []
if _top_v_pre and _top_v_pre[0][1] > 0:
    _pur_sum = sum(a for _, a in _top_v_pre[:10]) or 1
    _share = _top_v_pre[0][1] / _pur_sum * 100
    _about_lines.append(
        f"• Top vendor: {str(_top_v_pre[0][0])[:36]} "
        f"({_share:.1f}% of top-10 purchases)"
    )
add_text(slide, Inches(0.7), Inches(2.0), Inches(5.6), Inches(4.8),
         _about_lines, size=13, color=BLACK)

# Right column - Stats tiles
def stat_tile(x, y, w, h, value, label, color=NAVY):
    add_rect(slide, x, y, w, h, WHITE, line=GRAY_L)
    add_text(slide, x, y + Inches(0.2), w, Inches(0.65), value,
             size=24, bold=True, color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x, y + Inches(0.85), w, Inches(0.3), label,
             size=10, color=GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 6 tiles in 3x2 grid
tx = Inches(6.8); ty = Inches(1.5)
tw = Inches(2);   th = Inches(1.25)
gap = Inches(0.15)

stat_tile(tx,                     ty,                     tw, th, f"₹{to_cr(pla['Revenue']):.1f} Cr", "Annual Revenue (annualized)", NAVY)
stat_tile(tx + tw + gap,          ty,                     tw, th, f"{pl['Gross Margin %']:.1f}%",    "Gross Margin", GREEN)
stat_tile(tx + 2*(tw + gap),      ty,                     tw, th, f"{pl['EBITDA Margin %']:.1f}%",   "EBITDA Margin", NAVY)

# Stat tiles computed from the actual JSON — no hardcoded sample numbers.
# When a count is zero we show "—" (em dash) rather than "0+" which reads
# like an unfilled template.
def _count_label(n):
    return "—" if not n else str(int(n))

_n_vendors   = s.get('unique_vendor_count')   or len(s.get('top_vendors') or [])
_n_customers = s.get('unique_customer_count') or len(s.get('top_customers') or [])
_n_countries = len(s.get('country_rev') or {})

stat_tile(tx,                ty + th + gap, tw, th, _count_label(_n_vendors),   "Active Vendors",   NAVY)
stat_tile(tx + tw + gap,     ty + th + gap, tw, th, _count_label(_n_customers), "Active Customers", NAVY)
stat_tile(tx + 2*(tw + gap), ty + th + gap, tw, th, _count_label(_n_countries), "Countries",        GOLD)

# Operating Footprint computed from actual data — no hardcoded warehouse
# colors / product categories / departments. We show whatever the GL,
# sales and purchase exports actually tell us.
add_text(slide, Inches(6.8), Inches(4.4), Inches(6.3), Inches(0.4),
         "Operating Footprint", size=18, bold=True, color=NAVY)
add_rect(slide, Inches(6.8), Inches(4.9), Inches(6.3), Inches(2), BG_LIGHT)

_footprint = []
_countries = sorted(
    [c for c in (s.get('country_rev') or {}).keys() if c],
    key=lambda c: -(s.get('country_rev') or {}).get(c, 0),
)
if _countries:
    shown = ", ".join(_countries[:6])
    more = "" if len(_countries) <= 6 else f" + {len(_countries) - 6} more"
    _footprint.append(f"• Geographic reach: {shown}{more}")

_currencies = sorted(
    [c for c in (s.get('currency_rev') or {}).keys() if c],
    key=lambda c: -(s.get('currency_rev') or {}).get(c, 0),
)
if _currencies:
    _footprint.append(f"• Currencies billed in: {', '.join(_currencies[:8])}")

if _n_vendors:
    _footprint.append(f"• Active vendor relationships: {_n_vendors}")
if _n_customers:
    _footprint.append(f"• Active customer accounts: {_n_customers}")

if not _footprint:
    _footprint = ["• Operational data unavailable for this period."]

add_text(slide, Inches(6.95), Inches(5.0), Inches(6.1), Inches(1.9),
         _footprint, size=12, color=BLACK)

# =====================================================
# SLIDE 4 - FINANCIAL HIGHLIGHTS (KPI grid)
# =====================================================
slide = blank_slide()
add_title_bar(slide, f"Financial Highlights — {REPORTING_PERIOD}",
              "Headline metrics  |  11-month actuals & annualized run-rate")
add_footer(slide, 4)

# 6 large KPI cards
def kpi_card(x, y, w, h, label, value_act, value_ann, margin, accent=NAVY, neg=False):
    color = RED if neg else accent
    # Card background
    add_rect(slide, x, y, w, h, WHITE, line=GRAY_L)
    # Top accent bar
    add_rect(slide, x, y, w, Inches(0.12), color)
    # Label
    add_text(slide, x + Inches(0.15), y + Inches(0.2), w - Inches(0.3), Inches(0.35),
             label, size=12, bold=True, color=GRAY)
    # Main value
    add_text(slide, x + Inches(0.15), y + Inches(0.55), w - Inches(0.3), Inches(0.7),
             value_act, size=26, bold=True, color=color)
    # Annualized
    add_text(slide, x + Inches(0.15), y + Inches(1.3), w - Inches(0.3), Inches(0.3),
             f"Annualized: {value_ann}", size=10, color=GRAY, italic=True)
    # Margin
    if margin is not None:
        add_text(slide, x + Inches(0.15), y + Inches(1.6), w - Inches(0.3), Inches(0.3),
                 f"Margin: {margin}", size=11, bold=True, color=color)

cards = [
    ("REVENUE",         fmt_cr(pl['Revenue']),       fmt_cr(pla['Revenue']),       None,                            NAVY,  False),
    ("GROSS PROFIT",    fmt_cr(pl['Gross Profit']),  fmt_cr(pla['Gross Profit']),  fmt_pct(pl['Gross Margin %']),   GREEN, False),
    ("EBITDA",          fmt_cr(pl['EBITDA']),        fmt_cr(pla['EBITDA']),        fmt_pct(pl['EBITDA Margin %']),  NAVY,  pl['EBITDA'] < 0),
    ("EBIT",            fmt_cr(pl['EBIT']),          fmt_cr(pla['EBIT']),          fmt_pct(pl['EBIT Margin %']),    NAVY,  pl['EBIT'] < 0),
    ("PBT",             fmt_cr(pl['PBT']),           fmt_cr(pla['PBT']),           fmt_pct(pl['PBT Margin %']),     NAVY,  pl['PBT'] < 0),
    ("PAT (NET PROFIT)",fmt_cr(pl['PAT']),           fmt_cr(pla['PAT']),           fmt_pct(pl['PAT / Net Margin %']), GOLD,  pl['PAT'] < 0),
]
card_w = Inches(4.0); card_h = Inches(2.0)
x0 = Inches(0.5); y0 = Inches(1.4); gap = Inches(0.2)
for i, (lbl, va, vn, mg, color, neg) in enumerate(cards):
    row, col = divmod(i, 3)
    x = x0 + col * (card_w + gap)
    y = y0 + row * (card_h + gap)
    kpi_card(x, y, card_w, card_h, lbl, va, vn, mg, accent=color, neg=neg)

# Footnote
add_text(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.5),
         f"Note: 11-month actuals (Apr'24-Feb'25). Annualized = actuals × 12/11. {s['period']}",
         size=10, italic=True, color=GRAY)

# =====================================================
# SLIDE 5 - DETAILED P&L
# =====================================================
slide = blank_slide()
add_title_bar(slide, "Profit & Loss Statement",
              f"{REPORTING_PERIOD}  |  Figures in ₹ Crores  |  {s['period']}")
add_footer(slide, 5)

# P&L table
rows = [
    ("Revenue from Operations",          pl['Revenue'],       pla['Revenue'],       True,  False, True,  False),
    ("Less: Cost of Goods Sold (COGS)",  pl['COGS'],          pla['COGS'],          False, False, False, False),
    ("Gross Profit",                     pl['Gross Profit'],  pla['Gross Profit'],  True,  True,  True,  False),
    ("Less: Operating Expenses",         pl['Operating Expenses'], pla['Operating Expenses'], False, False, False, False),
    ("Add: Other Income",                pl['Other Income'],  pla['Other Income'],  False, False, False, False),
    ("Less: Other Expenses",             pl['Other Expenses'],pla['Other Expenses'],False, False, False, False),
    ("EBITDA",                           pl['EBITDA'],        pla['EBITDA'],        True,  True,  True,  False),
    ("Less: Depreciation & Amortization",pl['Depreciation'],  pla['Depreciation'],  False, False, False, False),
    ("EBIT (Operating Profit)",          pl['EBIT'],          pla['EBIT'],          True,  True,  True,  False),
    ("Less: Finance Costs",              pl['Finance Costs'], pla['Finance Costs'], False, False, False, False),
    ("Profit Before Tax (PBT)",          pl['PBT'],           pla['PBT'],           True,  True,  True,  False),
    ("Less: Tax (Indicative @25.17%)",   pl['Tax (Indicative @25.17%)'], pla['Tax (Indicative @25.17%)'], False, False, False, True),
    ("Profit After Tax (PAT)",           pl['PAT'],           pla['PAT'],           True,  True,  True,  False),
]

# Table dimensions
tx = Inches(0.5); ty = Inches(1.4)
col_w = [Inches(5.8), Inches(2.2), Inches(2.2), Inches(2.1)]
row_h = Inches(0.36)

# Header
header_y = ty
header_bg = NAVY
hx = tx
for i, (lbl, w) in enumerate(zip(["Particulars","11M Actual (₹ Cr)","Annualized (₹ Cr)","% of Revenue"], col_w)):
    add_rect(slide, hx, header_y, w, row_h, NAVY)
    add_text(slide, hx, header_y, w, row_h, lbl, size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER if i>0 else PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    hx += w

# Body rows
rev_act = pl['Revenue']
y = ty + row_h
for label, vact, vann, bold, subtotal, is_margin_row, italic in rows:
    fill = LIGHT_N if subtotal else (WHITE if (rows.index((label, vact, vann, bold, subtotal, is_margin_row, italic)) % 2 == 0) else BG_LIGHT)
    add_rect(slide, tx, y, sum(col_w, Emu(0)), row_h, fill, line=GRAY_L)
    add_text(slide, tx + Inches(0.15), y, col_w[0]-Inches(0.15), row_h, label,
             size=11, bold=bold, italic=italic, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)
    color = RED if vact < 0 else BLACK
    add_text(slide, tx + col_w[0], y, col_w[1], row_h,
             f"{to_cr(vact):,.2f}", size=11, bold=bold, color=color,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, tx + col_w[0] + col_w[1], y, col_w[2], row_h,
             f"{to_cr(vann):,.2f}", size=11, bold=bold, color=color,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    if is_margin_row and rev_act:
        pct = vact / rev_act * 100
        add_text(slide, tx + col_w[0] + col_w[1] + col_w[2], y, col_w[3], row_h,
                 fmt_pct(pct), size=11, bold=bold, color=color,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    y += row_h

# Right-side margin profile card
mx = Inches(10.9); my = Inches(1.4); mw = Inches(2.3)
add_rect(slide, mx, my, mw, Inches(4.7), BG_LIGHT)
add_text(slide, mx + Inches(0.15), my + Inches(0.1), mw - Inches(0.3), Inches(0.4),
         "Margin Profile", size=14, bold=True, color=NAVY)
mh = my + Inches(0.55)
margins = [
    ("Gross Margin", pl['Gross Margin %'], GREEN),
    ("EBITDA Margin", pl['EBITDA Margin %'], NAVY),
    ("EBIT Margin", pl['EBIT Margin %'], NAVY),
    ("PBT Margin", pl['PBT Margin %'], NAVY),
    ("PAT / Net Margin", pl['PAT / Net Margin %'], GOLD),
]
for lbl, val, col in margins:
    add_text(slide, mx + Inches(0.15), mh, mw - Inches(0.3), Inches(0.25),
             lbl, size=10, color=GRAY)
    add_text(slide, mx + Inches(0.15), mh + Inches(0.22), mw - Inches(0.3), Inches(0.45),
             fmt_pct(val), size=18, bold=True, color=(RED if val<0 else col))
    mh += Inches(0.85)

# =====================================================
# SLIDE 6 - MONTHLY REVENUE TREND (chart)
# =====================================================
slide = blank_slide()
add_title_bar(slide, "Revenue Trend & Momentum",
              f"Monthly revenue trajectory — {REPORTING_PERIOD}")
add_footer(slide, 6)

month_labels = {'2024-04':"Apr-24",'2024-05':"May-24",'2024-06':"Jun-24",'2024-07':"Jul-24",
                '2024-08':"Aug-24",'2024-09':"Sep-24",'2024-10':"Oct-24",'2024-11':"Nov-24",
                '2024-12':"Dec-24",'2025-01':"Jan-25",'2025-02':"Feb-25"}
mr_items = sorted(s['monthly_revenue'].items())
cats = [month_labels[ym] for ym, _ in mr_items]
vals = [round(to_cr(v), 2) for _, v in mr_items]

cdata = CategoryChartData()
cdata.categories = cats
cdata.add_series('Revenue (₹ Cr)', vals)

chart_shape = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.5), Inches(1.4), Inches(9.5), Inches(5.0), cdata
)
chart = chart_shape.chart
chart.has_title = False
chart.has_legend = False
plot = chart.plots[0]
plot.has_data_labels = True
plot.data_labels.font.size = Pt(10)
plot.data_labels.font.bold = True
plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
# bar colors
for series in chart.series:
    fill = series.format.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

# Insight panel right
ix = Inches(10.2); iy = Inches(1.4); iw = Inches(3.0)
add_rect(slide, ix, iy, iw, Inches(5.0), BG_LIGHT)
add_text(slide, ix + Inches(0.15), iy + Inches(0.15), iw - Inches(0.3), Inches(0.5),
         "Key Insights", size=16, bold=True, color=NAVY)

# Compute growth metrics
apr_rev = mr_items[0][1]
peak_rev = max(v for _,v in mr_items if v > 1000)
jan_rev = next((v for ym,v in mr_items if ym=='2025-01'), peak_rev)
growth = ((jan_rev / apr_rev) - 1) * 100 if apr_rev else 0

insight_y = iy + Inches(0.7)
insights = [
    (f"+{growth:.0f}%", "Apr'24 → Jan'25 growth", GREEN),
    (f"₹{to_cr(peak_rev):.2f} Cr", "Peak monthly revenue", NAVY),
    (f"₹{to_cr(sum(v for _,v in mr_items)/len([1 for _,v in mr_items if v>1000])):.2f} Cr", "Avg monthly run-rate", NAVY),
    ("H2 > H1", "Second-half momentum", GOLD),
]
for val, lbl, col in insights:
    add_text(slide, ix + Inches(0.15), insight_y, iw - Inches(0.3), Inches(0.5),
             val, size=22, bold=True, color=col)
    add_text(slide, ix + Inches(0.15), insight_y + Inches(0.45), iw - Inches(0.3), Inches(0.3),
             lbl, size=10, color=GRAY)
    insight_y += Inches(1.0)

# =====================================================
# SLIDE 7 - TOP CUSTOMERS
# =====================================================
slide = blank_slide()
add_title_bar(slide, "Top Customers", "Customer concentration based on recent invoice sample")
add_footer(slide, 7)

top_c = [(n, a) for n, a in (s.get('top_customers') or [])[:8] if a]
total_c = sum(a for _, a in top_c) or 1

# Left: table
tx = Inches(0.5); ty = Inches(1.4)
add_rect(slide, tx, ty, Inches(0.5), Inches(0.4), NAVY)
add_text(slide, tx, ty, Inches(0.5), Inches(0.4), "#", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(slide, tx + Inches(0.5), ty, Inches(4.5), Inches(0.4), NAVY)
add_text(slide, tx + Inches(0.5), ty, Inches(4.5), Inches(0.4), "Customer", size=11, bold=True, color=WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
add_rect(slide, tx + Inches(5), ty, Inches(1.7), Inches(0.4), NAVY)
add_text(slide, tx + Inches(5), ty, Inches(1.7), Inches(0.4), "Value", size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
add_rect(slide, tx + Inches(6.7), ty, Inches(1), Inches(0.4), NAVY)
add_text(slide, tx + Inches(6.7), ty, Inches(1), Inches(0.4), "Share", size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

if top_c:
    y = ty + Inches(0.4)
    for i, (name, amt) in enumerate(top_c, 1):
        bg = WHITE if i % 2 else BG_LIGHT
        add_rect(slide, tx, y, Inches(7.7), Inches(0.42), bg, line=GRAY_L)
        add_text(slide, tx, y, Inches(0.5), Inches(0.42), str(i), size=11, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, tx + Inches(0.55), y, Inches(4.45), Inches(0.42),
                 str(name)[:40], size=11, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, tx + Inches(5), y, Inches(1.65), Inches(0.42),
                 f"{amt:,.0f}", size=11, bold=True, color=BLACK, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        pct = amt/total_c*100
        add_text(slide, tx + Inches(6.7), y, Inches(0.95), Inches(0.42),
                 f"{pct:.1f}%", size=11, color=GOLD, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.42)

    # Right: pie chart (only when we have real categories to plot).
    cdata = CategoryChartData()
    cdata.categories = [str(n)[:25] for n, _ in top_c]
    cdata.add_series('Customers', [a for _, a in top_c])
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, Inches(8.5), Inches(1.4), Inches(4.8), Inches(4.6), cdata
    )
    chart = chart_shape.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Customer Mix"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = NAVY
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(8)
else:
    add_rect(slide, tx, ty + Inches(0.4), Inches(12.3), Inches(1.2), BG_LIGHT, line=GRAY_L)
    add_text(slide, tx + Inches(0.3), ty + Inches(0.7), Inches(12.0), Inches(0.6),
             "Insufficient sales data to compute top customers — re-upload your "
             "sales register with customer name + amount columns.",
             size=13, italic=True, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)

add_text(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
         "Note: Based on invoice line totals from sample period (Nov 2024 - Feb 2025). Currencies are as billed.",
         size=10, italic=True, color=GRAY)

# =====================================================
# SLIDE 8 - TOP VENDORS
# =====================================================
slide = blank_slide()
add_title_bar(slide, "Top Vendors / Suppliers", "Procurement concentration analysis")
add_footer(slide, 8)

top_v = [(n, a) for n, a in (s.get('top_vendors') or [])[:8] if a]
total_v = sum(a for _, a in top_v) or 1

tx = Inches(0.5); ty = Inches(1.4)
add_rect(slide, tx, ty, Inches(0.5), Inches(0.4), NAVY)
add_text(slide, tx, ty, Inches(0.5), Inches(0.4), "#", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(slide, tx + Inches(0.5), ty, Inches(4.5), Inches(0.4), NAVY)
add_text(slide, tx + Inches(0.5), ty, Inches(4.5), Inches(0.4), "Vendor", size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
add_rect(slide, tx + Inches(5), ty, Inches(1.7), Inches(0.4), NAVY)
add_text(slide, tx + Inches(5), ty, Inches(1.7), Inches(0.4), "Purchase Value", size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
add_rect(slide, tx + Inches(6.7), ty, Inches(1), Inches(0.4), NAVY)
add_text(slide, tx + Inches(6.7), ty, Inches(1), Inches(0.4), "Share", size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

if top_v:
    y = ty + Inches(0.4)
    for i, (name, amt) in enumerate(top_v, 1):
        bg = WHITE if i % 2 else BG_LIGHT
        add_rect(slide, tx, y, Inches(7.7), Inches(0.42), bg, line=GRAY_L)
        add_text(slide, tx, y, Inches(0.5), Inches(0.42), str(i), size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, tx + Inches(0.55), y, Inches(4.45), Inches(0.42), str(name)[:40], size=11, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, tx + Inches(5), y, Inches(1.65), Inches(0.42), f"{amt:,.0f}", size=11, bold=True, color=BLACK, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, tx + Inches(6.7), y, Inches(0.95), Inches(0.42), f"{amt/total_v*100:.1f}%", size=11, color=GOLD, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.42)
else:
    add_rect(slide, tx, ty + Inches(0.4), Inches(12.3), Inches(1.2), BG_LIGHT, line=GRAY_L)
    add_text(slide, tx + Inches(0.3), ty + Inches(0.7), Inches(12.0), Inches(0.6),
             "Insufficient purchase data to compute top vendors — re-upload your "
             "purchase register with vendor name + amount columns.",
             size=13, italic=True, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)

# Bar chart (only when we have real categories to plot — python-pptx
# raises "chart data contains no categories" otherwise).
if top_v:
    cdata = CategoryChartData()
    cdata.categories = [str(n)[:20] for n, _ in top_v]
    cdata.add_series('Purchase Value', [a for _, a in top_v])
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(8.5), Inches(1.4), Inches(4.8), Inches(4.6), cdata
    )
    chart = chart_shape.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Vendor Spend"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = NAVY
    chart.has_legend = False
    for series in chart.series:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = GOLD

add_text(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
         f"Note: Based on purchase invoice line totals ({REPORTING_PERIOD}).",
         size=10, italic=True, color=GRAY)

# =====================================================
# SLIDE 9 - GEOGRAPHY & CURRENCY
# =====================================================
slide = blank_slide()
add_title_bar(slide, "Geographic & Currency Mix", "Revenue distribution by country and billing currency")
add_footer(slide, 9)

# Country chart - left. Filter zero-amount countries so the pie chart
# doesn't render a slice of nothing, and skip entirely if the result
# is empty (otherwise python-pptx raises "chart data contains no
# categories").
country_data = [
    (c, a) for c, a in sorted(s['country_rev'].items(), key=lambda x: -x[1])[:10]
    if a
]
if country_data:
    cdata = CategoryChartData()
    cdata.categories = [c or 'IN' for c, _ in country_data]
    cdata.add_series('Revenue', [a for _, a in country_data])
    ch1 = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(0.5), Inches(1.4), Inches(6), Inches(5), cdata).chart
    ch1.has_title = True
    ch1.chart_title.text_frame.text = "Revenue by Country"
    ch1.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
    ch1.chart_title.text_frame.paragraphs[0].font.bold = True
    ch1.chart_title.text_frame.paragraphs[0].font.color.rgb = NAVY
    ch1.has_legend = True
    ch1.legend.position = XL_LEGEND_POSITION.RIGHT
    ch1.legend.font.size = Pt(10)
    ch1.plots[0].has_data_labels = True
    ch1.plots[0].data_labels.show_percentage = True
    ch1.plots[0].data_labels.font.size = Pt(9)
    ch1.plots[0].data_labels.font.bold = True
else:
    add_text(slide, Inches(0.5), Inches(3), Inches(6), Inches(0.6),
             "No country-level revenue available.",
             size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

# Currency chart - right
curr_data = [
    (c, a) for c, a in sorted(s['currency_rev'].items(), key=lambda x: -x[1])
    if a
]
if curr_data:
    cdata2 = CategoryChartData()
    cdata2.categories = [c or 'INR' for c, _ in curr_data]
    cdata2.add_series('Revenue', [a for _, a in curr_data])
    ch2 = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(7.0), Inches(1.4), Inches(6), Inches(5), cdata2).chart
    ch2.has_title = True
    ch2.chart_title.text_frame.text = "Revenue by Currency"
    ch2.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
    ch2.chart_title.text_frame.paragraphs[0].font.bold = True
    ch2.chart_title.text_frame.paragraphs[0].font.color.rgb = NAVY
    ch2.has_legend = True
    ch2.legend.position = XL_LEGEND_POSITION.RIGHT
    ch2.legend.font.size = Pt(10)
    ch2.plots[0].has_data_labels = True
    ch2.plots[0].data_labels.show_percentage = True
    ch2.plots[0].data_labels.font.size = Pt(9)
    ch2.plots[0].data_labels.font.bold = True
else:
    add_text(slide, Inches(7.0), Inches(3), Inches(6), Inches(0.6),
             "No currency-level revenue available.",
             size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
         "Note: Multi-currency invoicing across 10+ countries. Indicates strong export exposure alongside domestic base.",
         size=10, italic=True, color=GRAY)

# =====================================================
# SLIDE 10 - WORKING CAPITAL
# =====================================================
slide = blank_slide()
add_title_bar(slide, "Working Capital & Liquidity",
              "Receivables, payables, and the cash conversion cycle")
add_footer(slide, 10)

ar = s['total_ar_open']
ap = s['total_ap_open']
nwc = ar - ap
ar_days = (ar / pla['Revenue']) * 365 if pla['Revenue'] else 0
ap_days = (ap / pla['COGS']) * 365 if pla['COGS'] else 0

# 3 big tiles top
def big_tile(x, y, w, h, val, sub, accent):
    add_rect(slide, x, y, w, h, WHITE, line=GRAY_L)
    add_rect(slide, x, y, w, Inches(0.15), accent)
    add_text(slide, x, y + Inches(0.3), w, Inches(0.85),
             val, size=30, bold=True, color=accent,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x, y + Inches(1.2), w, Inches(0.4),
             sub, size=12, color=GRAY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

big_tile(Inches(0.5), Inches(1.4), Inches(4.0), Inches(1.8), fmt_cr(ar), "Open Receivables (AR)", NAVY)
big_tile(Inches(4.7), Inches(1.4), Inches(4.0), Inches(1.8), fmt_cr(ap), "Open Payables (AP)", GOLD)
big_tile(Inches(8.9), Inches(1.4), Inches(4.0), Inches(1.8), fmt_cr(nwc), "Net Working Capital", GREEN if nwc>=0 else RED)

# Cash cycle row
big_tile(Inches(0.5), Inches(3.4), Inches(4.0), Inches(1.8), f"{ar_days:.0f} days", "AR Days (DSO)", NAVY)
big_tile(Inches(4.7), Inches(3.4), Inches(4.0), Inches(1.8), f"{ap_days:.0f} days", "AP Days (DPO)", GOLD)
ccc = ar_days - ap_days
big_tile(Inches(8.9), Inches(3.4), Inches(4.0), Inches(1.8), f"{ccc:.0f} days", "Cash Conversion (CCC)",
         GREEN if ccc<60 else (GOLD if ccc<90 else RED))

# Commentary
add_rect(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.5), BG_LIGHT)
add_text(slide, Inches(0.7), Inches(5.5), Inches(12), Inches(0.4),
         "Working Capital Commentary", size=14, bold=True, color=NAVY)
comm = [
    f"• Receivables of {fmt_cr(ar)} against payables of {fmt_cr(ap)} — net working capital deployment of {fmt_cr(nwc)}.",
    f"• DSO at {ar_days:.0f} days suggests {'tight collection discipline' if ar_days < 60 else 'opportunity to accelerate collections'}.",
    f"• Cash conversion cycle of {ccc:.0f} days — {'efficient' if ccc < 60 else 'requires attention'} working capital position.",
]
y = Inches(5.9)
for c in comm:
    add_text(slide, Inches(0.7), y, Inches(12), Inches(0.35), c, size=11, color=BLACK)
    y += Inches(0.32)

# =====================================================
# ANALYTICAL SLIDES (auto-generated)
# =====================================================

# ---- Balance Sheet snapshot ----
_bs = s.get('balance_sheet')
if _bs:
    slide = blank_slide()
    add_title_bar(slide, "Balance Sheet Snapshot",
                  "Closing balances from GL  •  ₹ Crores")
    add_footer(slide, 0)
    # Two columns: Assets | Liabilities & Equity
    cx_l = Inches(0.5); cy = Inches(1.4); cw = Inches(6.15); ch_h = Inches(4.6)
    cx_r = Inches(6.85)
    # Left card (Assets)
    add_rect(slide, cx_l, cy, cw, ch_h, BG_LIGHT)
    add_rect(slide, cx_l, cy, cw, Inches(0.4), NAVY)
    add_text(slide, cx_l, cy, cw, Inches(0.4), "ASSETS",
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rows_left = [
        ("Fixed Assets (Net)", _bs['fixedAssetsNet']),
        ("Inventory",          _bs['inventory']),
        ("Receivables",        _bs['receivables']),
        ("Cash & Bank",        _bs['cashAndBank']),
        ("Other Assets",       _bs['otherAssets']),
        ("Total Assets",       _bs['totalAssets']),
    ]
    ry = cy + Inches(0.55)
    for i, (lbl, val) in enumerate(rows_left):
        bold = (lbl == "Total Assets")
        add_text(slide, cx_l + Inches(0.3), ry, cw - Inches(2.5), Inches(0.4),
                 lbl, size=12, bold=bold, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, cx_l + cw - Inches(2.3), ry, Inches(2.0), Inches(0.4),
                 fmt_cr(val), size=12, bold=bold, color=BLACK,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        if bold:
            add_rect(slide, cx_l + Inches(0.3), ry, cw - Inches(0.6), Emu(12700), NAVY)
        ry += Inches(0.55)

    # Right card (Liabilities & Equity)
    add_rect(slide, cx_r, cy, cw, ch_h, BG_LIGHT)
    add_rect(slide, cx_r, cy, cw, Inches(0.4), GOLD)
    add_text(slide, cx_r, cy, cw, Inches(0.4), "LIABILITIES & EQUITY",
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rows_right = [
        ("Equity",                _bs['equity']),
        ("Long-Term Debt",        _bs['longTermDebt']),
        ("Trade Payables",        _bs['payables']),
        ("Other Current Liab.",   _bs['otherCurrentLiab']),
        ("Total",                 _bs['totalLiabilitiesAndEquity']),
    ]
    ry = cy + Inches(0.55)
    for lbl, val in rows_right:
        bold = (lbl == "Total")
        add_text(slide, cx_r + Inches(0.3), ry, cw - Inches(2.5), Inches(0.4),
                 lbl, size=12, bold=bold, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, cx_r + cw - Inches(2.3), ry, Inches(2.0), Inches(0.4),
                 fmt_cr(val), size=12, bold=bold, color=BLACK,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        ry += Inches(0.55)

    # Ratios footer
    add_text(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.4),
             f"Current Ratio: {_bs['currentRatio']:.2f}x   |   Debt/Equity: {_bs['debtToEquity']:.2f}x   |   Total Assets: {fmt_cr(_bs['totalAssets'])}",
             size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ---- Performance vs Benchmarks ----
_bm = s.get('benchmarks') or []
if _bm:
    slide = blank_slide()
    add_title_bar(slide, "Performance vs Industry Benchmarks",
                  "Where the business stands against typical SME / mid-market norms")
    add_footer(slide, 0)
    # Table
    tx = Inches(0.5); ty = Inches(1.4)
    col_w = [Inches(4.5), Inches(2.0), Inches(2.0), Inches(1.8), Inches(2.0)]
    row_h = Inches(0.42)
    hdrs = ["Metric", "Actual", "Benchmark", "Status", "Gap"]
    hx = tx
    for i, (lbl, w) in enumerate(zip(hdrs, col_w)):
        add_rect(slide, hx, ty, w, row_h, NAVY)
        add_text(slide, hx, ty, w, row_h, lbl, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        hx += w
    y = ty + row_h
    for idx, b in enumerate(_bm):
        bg = WHITE if idx % 2 == 0 else BG_LIGHT
        status = b.get('status', 'ok')
        status_color = GREEN if status == 'good' else (GOLD if status == 'ok' else RED)
        status_label = 'ON TARGET' if status == 'good' else ('CAUTION' if status == 'ok' else 'BELOW')
        x = tx
        for i, w in enumerate(col_w):
            add_rect(slide, x, y, w, row_h, bg, line=GRAY_L)
            x += w
        x = tx
        add_text(slide, x + Inches(0.15), y, col_w[0] - Inches(0.15), row_h,
                 str(b.get('metric', '')), size=11, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[0]
        add_text(slide, x, y, col_w[1], row_h,
                 str(b.get('actual', '')), size=11, bold=True, color=BLACK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[1]
        add_text(slide, x, y, col_w[2], row_h,
                 str(b.get('benchmark', '')), size=11, color=GRAY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[2]
        add_text(slide, x, y, col_w[3], row_h, status_label, size=10, bold=True, color=status_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[3]
        add_text(slide, x + Inches(0.1), y, col_w[4] - Inches(0.1), row_h,
                 str(b.get('gap', ''))[:35], size=10, color=BLACK,
                 anchor=MSO_ANCHOR.MIDDLE)
        y += row_h


# ---- Cost Structure ----
_cs = s.get('cost_structure') or []
if _cs:
    slide = blank_slide()
    add_title_bar(slide, "Cost Structure",
                  "Operating cost breakdown as % of revenue  •  Watch-list flagged in red")
    add_footer(slide, 0)
    tx = Inches(0.5); ty = Inches(1.4)
    col_w = [Inches(5.5), Inches(2.5), Inches(2.0), Inches(2.3)]
    row_h = Inches(0.38)
    hdrs = ["Category", "Amount (₹ Cr)", "% of Revenue", "Watchlist"]
    hx = tx
    for i, (lbl, w) in enumerate(zip(hdrs, col_w)):
        add_rect(slide, hx, ty, w, row_h, NAVY)
        add_text(slide, hx, ty, w, row_h, lbl, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        hx += w
    y = ty + row_h
    for idx, c in enumerate(_cs[:13]):
        bg = WHITE if idx % 2 == 0 else BG_LIGHT
        watch = c.get('isWatchlist')
        x = tx
        for w in col_w:
            add_rect(slide, x, y, w, row_h, bg, line=GRAY_L)
            x += w
        x = tx
        add_text(slide, x + Inches(0.15), y, col_w[0] - Inches(0.15), row_h,
                 str(c.get('category', '')), size=11, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[0]
        add_text(slide, x, y, col_w[1] - Inches(0.1), row_h,
                 fmt_cr(c.get('amount', 0)), size=11, color=BLACK,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[1]
        add_text(slide, x, y, col_w[2], row_h,
                 f"{c.get('percentOfRevenue', 0):.1f}%", size=11, bold=watch, color=RED if watch else BLACK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[2]
        add_text(slide, x, y, col_w[3], row_h,
                 "REVIEW" if watch else "OK",
                 size=10, bold=True, color=RED if watch else GREEN,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        y += row_h


# ---- Critical Issues ----
_ci = s.get('critical_issues') or []
if _ci:
    slide = blank_slide()
    add_title_bar(slide, f"Critical Issues Requiring Board Action ({len(_ci)})",
                  "Auto-detected from your data  •  Ranked by severity")
    add_footer(slide, 0)
    # Stack issues vertically; show top 3-4 in detail
    iy = Inches(1.4); ih = Inches(1.35); igap = Inches(0.1)
    for i, issue in enumerate(_ci[:4]):
        sev = issue.get('severity', 'medium')
        accent = RED if sev == 'high' else (GOLD if sev == 'medium' else NAVY)
        add_rect(slide, Inches(0.5), iy, Inches(12.3), ih, BG_LIGHT)
        add_rect(slide, Inches(0.5), iy, Inches(0.15), ih, accent)
        # Rank badge
        add_rect(slide, Inches(0.8), iy + Inches(0.2), Inches(0.5), Inches(0.5), accent)
        add_text(slide, Inches(0.8), iy + Inches(0.2), Inches(0.5), Inches(0.5),
                 str(issue.get('rank', i + 1)), size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Title + body
        add_text(slide, Inches(1.5), iy + Inches(0.1), Inches(11), Inches(0.35),
                 str(issue.get('title', '')), size=13, bold=True, color=NAVY)
        add_text(slide, Inches(1.5), iy + Inches(0.45), Inches(11), Inches(0.35),
                 f"Why: {issue.get('rootCause', '')}"[:160], size=10, color=BLACK)
        add_text(slide, Inches(1.5), iy + Inches(0.78), Inches(11), Inches(0.3),
                 f"Action: {issue.get('recommendedAction', '')}"[:160], size=10, italic=True, color=GRAY)
        add_text(slide, Inches(1.5), iy + Inches(1.05), Inches(11), Inches(0.3),
                 f"Impact: {issue.get('potentialImpact', '')}"[:160], size=10, bold=True, color=GREEN)
        iy += ih + igap


# ---- Growth Opportunities ----
_go = s.get('growth_opportunities') or []
if _go:
    slide = blank_slide()
    add_title_bar(slide, f"Growth Opportunities ({len(_go)})",
                  "Auto-detected upside paths from the data")
    add_footer(slide, 0)
    iy = Inches(1.4); ih = Inches(1.35); igap = Inches(0.1)
    for i, op in enumerate(_go[:4]):
        add_rect(slide, Inches(0.5), iy, Inches(12.3), ih, BG_LIGHT)
        add_rect(slide, Inches(0.5), iy, Inches(0.15), ih, GREEN)
        add_rect(slide, Inches(0.8), iy + Inches(0.2), Inches(0.5), Inches(0.5), GREEN)
        add_text(slide, Inches(0.8), iy + Inches(0.2), Inches(0.5), Inches(0.5),
                 str(op.get('rank', i + 1)), size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(1.5), iy + Inches(0.1), Inches(11), Inches(0.35),
                 str(op.get('title', '')), size=13, bold=True, color=NAVY)
        add_text(slide, Inches(1.5), iy + Inches(0.45), Inches(11), Inches(0.35),
                 f"Why: {op.get('rationale', '')}"[:160], size=10, color=BLACK)
        add_text(slide, Inches(1.5), iy + Inches(0.78), Inches(11), Inches(0.3),
                 f"Approach: {op.get('approach', '')}"[:160], size=10, italic=True, color=GRAY)
        add_text(slide, Inches(1.5), iy + Inches(1.05), Inches(11), Inches(0.3),
                 f"Upside: {op.get('potentialUpside', '')}"[:160], size=10, bold=True, color=GREEN)
        iy += ih + igap


# ---- Improvement Initiatives (quantified) ----
_imp = s.get('improvements') or []
if _imp:
    slide = blank_slide()
    total_savings = sum(m.get('savings', 0) for m in _imp)
    add_title_bar(slide, "Improvement Initiatives (Quantified)",
                  f"Total potential upside: {fmt_cr(total_savings)} per annum")
    add_footer(slide, 0)
    tx = Inches(0.5); ty = Inches(1.4)
    col_w = [Inches(5.5), Inches(2.5), Inches(2.0), Inches(2.3)]
    row_h = Inches(0.42)
    hdrs = ["Action", "Impact (₹ Cr p.a.)", "Timeline", "Difficulty"]
    hx = tx
    for i, (lbl, w) in enumerate(zip(hdrs, col_w)):
        add_rect(slide, hx, ty, w, row_h, NAVY)
        add_text(slide, hx, ty, w, row_h, lbl, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        hx += w
    y = ty + row_h
    for idx, m in enumerate(_imp[:10]):
        bg = WHITE if idx % 2 == 0 else BG_LIGHT
        diff = m.get('difficulty', 'medium')
        diff_color = GREEN if diff == 'easy' else (GOLD if diff == 'medium' else RED)
        x = tx
        for w in col_w:
            add_rect(slide, x, y, w, row_h, bg, line=GRAY_L)
            x += w
        x = tx
        add_text(slide, x + Inches(0.15), y, col_w[0] - Inches(0.15), row_h,
                 str(m.get('action', '')), size=11, bold=True, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[0]
        add_text(slide, x, y, col_w[1] - Inches(0.1), row_h,
                 fmt_cr(m.get('savings', 0)), size=11, bold=True, color=GREEN,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[1]
        add_text(slide, x, y, col_w[2], row_h,
                 str(m.get('timeline', '')), size=11, color=BLACK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[2]
        add_text(slide, x, y, col_w[3], row_h,
                 diff.upper(), size=10, bold=True, color=diff_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        y += row_h


# ---- AI Summary commentary slide ----
_ai = s.get('ai_summary') or ''
if _ai:
    slide = blank_slide()
    add_title_bar(slide, "AI Summary", "Automated read of your financials and analysis")
    add_footer(slide, 0)
    # Big centered text block
    add_rect(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5), BG_LIGHT)
    paragraphs = _ai.split("\n\n")
    y = Inches(1.7); ph = Inches(1.4)
    for i, p in enumerate(paragraphs[:3]):
        add_text(slide, Inches(1), y, Inches(11.3), ph, p[:500], size=14, color=BLACK)
        y += ph + Inches(0.2)


# =====================================================
# OPTIONAL AGING + CASH FLOW slides
# =====================================================
_ca = s.get('customer_aging')
_va = s.get('vendor_aging')
_cf = s.get('cash_flow')

def _aging_slide(title_text, subtitle_text, aging, accent_hex, party_label):
    """Render one aging slide with bucket bar chart + top parties table."""
    if not aging:
        return
    slide = blank_slide()
    add_title_bar(slide, title_text, subtitle_text)
    add_footer(slide, 0)

    # KPI strip across top
    total = aging.get('totalOutstanding', 0)
    pc = aging.get('partyCount', 0)
    asof = aging.get('asOfDate') or ''
    buckets = aging.get('buckets', {})
    overdue = sum(buckets.get(b, 0) for b in ('31-60', '61-90', '91-180', '180+'))
    pct_overdue = (overdue / total * 100) if total else 0

    tx = Inches(0.5); ty = Inches(1.4); tw = Inches(2.95); th = Inches(1.0)
    gap = Inches(0.1)
    tiles = [
        ("Total Outstanding", fmt_cr(total), accent_hex),
        ("Parties", str(pc), NAVY),
        ("31+ Days Overdue", fmt_cr(overdue), RED),
        ("% Overdue", f"{pct_overdue:.1f}%", RED if pct_overdue > 30 else GOLD),
    ]
    for i, (lbl, val, c) in enumerate(tiles):
        x = tx + i * (tw + gap)
        add_rect(slide, x, ty, tw, th, WHITE, line=GRAY_L)
        add_rect(slide, x, ty, tw, Inches(0.1), c)
        add_text(slide, x, ty + Inches(0.15), tw, Inches(0.32), lbl, size=9, color=GRAY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x, ty + Inches(0.45), tw, Inches(0.5), val, size=18, bold=True, color=c,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Bucket chart (clustered column)
    bucket_order = ['Current', '0-30', '31-60', '61-90', '91-180', '180+']
    cd = CategoryChartData()
    cd.categories = bucket_order
    cd.add_series('Outstanding (₹ Cr)', [round(to_cr(buckets.get(b, 0)), 2) for b in bucket_order])
    ch_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(2.55), Inches(7.5), Inches(3.6), cd
    )
    ch = ch_shape.chart
    ch.has_title = False
    ch.has_legend = False
    ch.plots[0].has_data_labels = True
    ch.plots[0].data_labels.font.size = Pt(10)
    ch.plots[0].data_labels.font.bold = True
    ch.plots[0].data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    for series in ch.series:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = accent_hex

    # Top parties table
    tpx = Inches(8.2); tpy = Inches(2.55); tpw = Inches(4.7)
    add_text(slide, tpx, tpy, tpw, Inches(0.35),
             f"Top {party_label}s by Outstanding", size=12, bold=True, color=NAVY)
    row_y = tpy + Inches(0.4)
    top = (aging.get('topParties') or [])[:8]
    for idx, p in enumerate(top, 1):
        bg = BG_LIGHT if idx % 2 == 0 else WHITE
        add_rect(slide, tpx, row_y, tpw, Inches(0.4), bg, line=GRAY_L)
        add_text(slide, tpx + Inches(0.1), row_y, Inches(0.3), Inches(0.4),
                 str(idx), size=10, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, tpx + Inches(0.45), row_y, Inches(2.55), Inches(0.4),
                 str(p.get('name', ''))[:30], size=10, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)
        bk = p.get('bucket', 'Current')
        bk_color = GREEN if bk == 'Current' else (GOLD if bk in ('0-30', '31-60') else RED)
        add_text(slide, tpx + Inches(3.0), row_y, Inches(0.7), Inches(0.4),
                 bk, size=9, bold=True, color=bk_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, tpx + Inches(3.7), row_y, Inches(0.95), Inches(0.4),
                 fmt_cr(p.get('amount', 0)), size=10, bold=True, color=BLACK,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        row_y += Inches(0.4)

    if asof:
        add_text(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.4),
                 f"As of {asof}  •  Buckets in days past due. Current = not yet due.",
                 size=10, italic=True, color=GRAY)


_aging_slide("Customer Aging — Receivables",
             "Aged outstanding receivables and top overdue customers",
             _ca, NAVY, "Customer")
_aging_slide("Vendor Aging — Payables",
             "Aged outstanding payables and top overdue vendors",
             _va, GOLD, "Vendor")

# Cash Flow Projection slide (only if both)
if _cf and _ca and _va:
    slide = blank_slide()
    add_title_bar(slide, "Cash-Flow Projection (Next 90 Days)",
                  "Expected collections vs payments — net cash position by 30-day window")
    add_footer(slide, 0)

    # 3 KPI tiles
    tx = Inches(0.5); ty = Inches(1.4); tw = Inches(4.0); th = Inches(1.4)
    gap = Inches(0.15)
    total_c = _cf.get('totalCollections', 0)
    total_p = _cf.get('totalPayments', 0)
    net = _cf.get('netCashFlow', 0)
    tiles = [
        ("Expected Collections (90d)", fmt_cr(total_c), GREEN),
        ("Expected Payments (90d)", fmt_cr(total_p), RED),
        ("Net Cash Flow", fmt_cr(net), GREEN if net >= 0 else RED),
    ]
    for i, (lbl, val, c) in enumerate(tiles):
        x = tx + i * (tw + gap)
        add_rect(slide, x, ty, tw, th, WHITE, line=GRAY_L)
        add_rect(slide, x, ty, tw, Inches(0.12), c)
        add_text(slide, x, ty + Inches(0.25), tw, Inches(0.4), lbl, size=11, color=GRAY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x, ty + Inches(0.65), tw, Inches(0.65), val, size=22, bold=True, color=c,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Chart: clustered columns Collections / Payments / Net per window
    wins = _cf.get('windows', [])
    cd = CategoryChartData()
    cd.categories = [w['label'] for w in wins]
    cd.add_series('Collections', [round(to_cr(w['collections']), 2) for w in wins])
    cd.add_series('Payments', [round(to_cr(w['payments']), 2) for w in wins])
    cd.add_series('Net', [round(to_cr(w['net']), 2) for w in wins])
    ch_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(3.0), Inches(12.3), Inches(3.4), cd
    )
    ch = ch_shape.chart
    ch.has_title = False
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.font.size = Pt(11)
    ch.plots[0].has_data_labels = True
    ch.plots[0].data_labels.font.size = Pt(10)
    ch.plots[0].data_labels.font.bold = True
    series_colors = [GREEN, RED, NAVY]
    for series, color in zip(ch.series, series_colors):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color

    add_text(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
             "Heuristic: Current + 0-30 buckets convert to cash within 30 days. Older buckets recover at 70% / 50% rates.",
             size=10, italic=True, color=GRAY)


# =====================================================
# OPTIONAL BUDGET vs ACTUAL slide (only when budget uploaded)
# =====================================================
_variance = s.get('variance') or []
if _variance:
    slide = blank_slide()
    add_title_bar(slide, "Budget vs Actual — Variance Analysis",
                  "Computed against annualized actuals  |  Favorable = better than budget")
    add_footer(slide, 11)

    # Table layout
    tx = Inches(0.5); ty = Inches(1.4)
    col_w = [Inches(4.0), Inches(1.2), Inches(1.7), Inches(1.7), Inches(1.7), Inches(1.2), Inches(1.4)]
    row_h = Inches(0.32)

    # Header row
    hx = tx
    hdrs = ["Line Item", "Type", "Budget (₹ Cr)", "Actual (₹ Cr)", "Variance (₹ Cr)", "Var %", "Status"]
    for i, (lbl, w) in enumerate(zip(hdrs, col_w)):
        add_rect(slide, hx, ty, w, row_h, NAVY)
        add_text(slide, hx, ty, w, row_h, lbl, size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        hx += w

    # Body rows (top 14)
    y = ty + row_h
    for idx, v in enumerate(_variance[:14]):
        fill = WHITE if idx % 2 == 0 else BG_LIGHT
        favorable = v['favorable']
        color_var = GREEN if favorable else RED
        x = tx
        # Line item
        add_rect(slide, x, y, col_w[0], row_h, fill, line=GRAY_L)
        add_text(slide, x + Inches(0.1), y, col_w[0] - Inches(0.1), row_h,
                 str(v['lineItem'])[:48], size=10, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[0]
        # Type
        add_rect(slide, x, y, col_w[1], row_h, fill, line=GRAY_L)
        add_text(slide, x, y, col_w[1], row_h,
                 v['kind'].upper(), size=9, bold=True,
                 color=GREEN if v['kind'] == 'income' else GOLD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[1]
        # Budget
        add_rect(slide, x, y, col_w[2], row_h, fill, line=GRAY_L)
        add_text(slide, x, y, col_w[2] - Inches(0.1), row_h,
                 f"{to_cr(v['budget']):,.2f}", size=10, color=BLACK,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[2]
        # Actual
        add_rect(slide, x, y, col_w[3], row_h, fill, line=GRAY_L)
        add_text(slide, x, y, col_w[3] - Inches(0.1), row_h,
                 f"{to_cr(v['actual']):,.2f}", size=10, color=BLACK,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[3]
        # Variance
        add_rect(slide, x, y, col_w[4], row_h, fill, line=GRAY_L)
        add_text(slide, x, y, col_w[4] - Inches(0.1), row_h,
                 f"{to_cr(v['variance']):,.2f}", size=10, bold=True, color=color_var,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[4]
        # Var %
        add_rect(slide, x, y, col_w[5], row_h, fill, line=GRAY_L)
        sign = '+' if v['variancePct'] >= 0 else ''
        add_text(slide, x, y, col_w[5] - Inches(0.1), row_h,
                 f"{sign}{v['variancePct']:.1f}%", size=10, color=color_var,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[5]
        # Status
        add_rect(slide, x, y, col_w[6], row_h, fill, line=GRAY_L)
        add_text(slide, x, y, col_w[6], row_h,
                 "FAVORABLE" if favorable else "ADVERSE",
                 size=9, bold=True, color=color_var,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        y += row_h

    # Footnote
    add_text(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
             f"Showing top {min(len(_variance), 14)} variance lines (sorted by absolute variance). Full variance available in MIS workbook.",
             size=10, italic=True, color=GRAY)

# =====================================================
# SLIDE - GROWTH STORY / ACHIEVEMENTS
# =====================================================
slide = blank_slide()
add_title_bar(slide, f"{REPORTING_PERIOD} — Key Achievements", "What we accomplished this year")
add_footer(slide, 11)

# 2x2 achievement boxes
def ach_box(x, y, w, h, icon_label, title, body, color):
    add_rect(slide, x, y, w, h, BG_LIGHT)
    add_rect(slide, x, y, w, Inches(0.1), color)
    add_text(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.8), Inches(0.5),
             icon_label, size=22, bold=True, color=color)
    add_text(slide, x + Inches(1.2), y + Inches(0.25), w - Inches(1.4), Inches(0.45),
             title, size=15, bold=True, color=NAVY)
    add_text(slide, x + Inches(0.3), y + Inches(0.85), w - Inches(0.6), h - Inches(0.95),
             body, size=12, color=BLACK)

bx = Inches(0.5); by = Inches(1.4); bw = Inches(6.15); bh = Inches(2.5)
gap_x = Inches(0.2); gap_y = Inches(0.2)

# Pull 4 user-editable achievements from OUTLOOK (fall back to safe defaults)
_acks = (OUTLOOK.get('achievements') or [])
while len(_acks) < 4:
    _acks.append({"title": "", "body": ""})
_ack_colors = [NAVY, GREEN, GOLD, NAVY2]
_positions = [
    (bx, by),
    (bx + bw + gap_x, by),
    (bx, by + bh + gap_y),
    (bx + bw + gap_x, by + bh + gap_y),
]
for i, ((x, y), color) in enumerate(zip(_positions, _ack_colors)):
    ach_box(x, y, bw, bh, str(i + 1),
            _acks[i].get("title", "") or f"Achievement {i + 1}",
            _acks[i].get("body", "") or "—",
            color)

# =====================================================
# SLIDE 12 - OUTLOOK FY 25-26
# =====================================================
slide = blank_slide()
add_title_bar(slide, "Strategic Outlook — Next Year",
              "Where we are headed")
add_footer(slide, 12)

# Three pillar columns - sourced from OUTLOOK config
_g = OUTLOOK.get('growth', {})
_p = OUTLOOK.get('profitability', {})
_c = OUTLOOK.get('capability', {})
pillars = [
    (_g.get('tag', 'GROWTH'), _g.get('title', 'Revenue & Market Expansion'), _g.get('bullets', []) or [''], NAVY),
    (_p.get('tag', 'PROFITABILITY'), _p.get('title', 'Margin & Cost Discipline'), _p.get('bullets', []) or [''], GREEN),
    (_c.get('tag', 'CAPABILITY'), _c.get('title', 'People, Tech & Infra'), _c.get('bullets', []) or [''], GOLD),
]
px = Inches(0.5); py = Inches(1.5); pw = Inches(4.15); ph = Inches(5.3)
for i, (tag, title, bullets, color) in enumerate(pillars):
    x = px + i * (pw + Inches(0.2))
    add_rect(slide, x, py, pw, ph, BG_LIGHT)
    add_rect(slide, x, py, pw, Inches(0.5), color)
    add_text(slide, x, py, pw, Inches(0.5), tag, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(0.2), py + Inches(0.65), pw - Inches(0.4), Inches(0.5),
             title, size=14, bold=True, color=NAVY)
    body = [f"• {b}" for b in bullets]
    add_text(slide, x + Inches(0.2), py + Inches(1.15), pw - Inches(0.4), Inches(3.2),
             body, size=12, color=BLACK)

# =====================================================
# SLIDE 13 - RISKS, ASKS, DISCUSSION
# =====================================================
slide = blank_slide()
add_title_bar(slide, "Risks, Asks & Board Discussion",
              "Where we need the board's input")
add_footer(slide, 13)

# Two columns: Risks (left) and Asks (right)
add_rect(slide, Inches(0.5), Inches(1.4), Inches(6.15), Inches(5.5), BG_LIGHT)
add_text(slide, Inches(0.7), Inches(1.55), Inches(6), Inches(0.5),
         "⚠ KEY RISKS", size=18, bold=True, color=RED)

risks = [f"• {r}" for r in (OUTLOOK.get('risks') or [])]
y = Inches(2.1)
for r in risks:
    add_text(slide, Inches(0.7), y, Inches(5.9), Inches(0.55), r, size=12, color=BLACK)
    y += Inches(0.7)

add_rect(slide, Inches(6.85), Inches(1.4), Inches(6.0), Inches(5.5), BG_LIGHT)
add_text(slide, Inches(7.05), Inches(1.55), Inches(5.8), Inches(0.5),
         "🤝 BOARD ASKS", size=18, bold=True, color=NAVY)

asks = [f"• {a}" for a in (OUTLOOK.get('asks') or [])]
y = Inches(2.1)
for a in asks:
    add_text(slide, Inches(7.05), y, Inches(5.7), Inches(0.55), a, size=12, color=BLACK)
    y += Inches(0.7)

# =====================================================
# SLIDE 14 - THANK YOU
# =====================================================
slide = blank_slide()
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(slide, 0, Inches(4.2), SLIDE_W, Inches(0.08), GOLD)

add_text(slide, Inches(0.7), Inches(2.5), Inches(12), Inches(1.5),
         "THANK YOU", size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.7), Inches(4.5), Inches(12), Inches(0.6),
         "Questions & Discussion", size=22, color=GOLD_L, align=PP_ALIGN.CENTER, italic=True)
add_text(slide, Inches(0.7), Inches(5.5), Inches(12), Inches(0.5),
         f"{COMPANY_NAME}  |  {PREPARED_FOR}  |  {REPORTING_PERIOD} Review",
         size=14, color=LIGHT_N, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.7), Inches(6.1), Inches(12), Inches(0.5),
         f"Prepared: {s['period']}", size=11, italic=True, color=GOLD_L, align=PP_ALIGN.CENTER)

# Save
prs.save(OUT)
print(f"Deck saved: {OUT}")
print(f"Total slides: {len(prs.slides)}")
